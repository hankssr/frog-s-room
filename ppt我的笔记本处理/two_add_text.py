# """
# 这个文件可以处理“我的笔记本.ptt' 文件以下功能
#     1、月计划页：复制11页，然后在每一页特定位置添加月份文本框，添加当月日历
#     2、每日页:根据从当天到2022年12月31日复制日计划页，并在每页特定位置标记？？？年？月？日星期几。
#     3、在每日页特定位置，打印名人名言。单独从名人名言库中读取。
#     4、月日历链接到对应日计划页
# """

import calendar
import datetime
import random
from pptx.enum.dml import MSO_THEME_COLOR
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Pt,Cm
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_PARAGRAPH_ALIGNMENT
# 打开已存在ppt
ppt = Presentation('2笔记本_增加页面.pptx')
print(f'这个ppt共有{len(ppt.slides)}页,0为第一页')
'''
每个月至少有28天，所以初始值设成这个月的28号，然后尝试对这个日期加1天，加2天，加3天，什么时候日期跳到下个月了就停止，这样就知道这个月的天数了。
判断一个月有多少天
'''

# 全局变量和常数定义
week_list = ["星期一", "星期二", "星期三", "星期四", "星期五", "星期六", "星期日"]
year = 2022
num = 0
My_time = []  # 生成[2022年5月1日星期日]格式的My_time列表;
Sour_time = []  # 生成[2022-5-1]格式的My_time列表;
days_plan = 13  # 日计划起始页,第14页,还有一个位置有同名参数,记得同时改
List_Month = ['五月', '六月', '七月', '八月', '九月', '十月', '十一月', '十二月']  # 定义月份字符数组
月计划起始页 = 5  # 月计划起始页6页,还有一个位置有同名参数,记得同时改
月份 = 5  # 从5月份开始计算
习惯养成开始页 = 258
菜谱开始页 = 266
股票开始页 = 306
康奈尔笔记开始页 = 336
方格纸开始页 = 366
米字格开始页 = 371
横线开始页=376
摘录开始页=381


# 判断一个月有多少天
def month_days(year, month):
    begin = datetime.date(year, month, 28)
    for i in range(1, 5):
        end = begin + datetime.timedelta(days=i)
        if end.month > month:
            break
    result = 28 + i - 1
    return result  # 返回一个月的天数int


# 将年月日星期数据存入My_time列表
for i in range(5, 13):  # 5-12月
    for n in range(1, month_days(year, i) + 1):  # 每个月天数
        week = week_list[datetime.date(year, i, n).weekday()]
        My_time.append(str(year) + '年' + str(i) + '月' + str(n) + '日' + week)
        Sour_time.append(str(year) + '-' + str(i) + '-' + str(n))
print(f'年月日格式的日期数组My_time是{My_time}')
print(f'日期格式的年月日数组Sour_time是{Sour_time}')

#在第一页增加股票,摘抄等链接

def first_Page_hypro(left, top, width, height,text,page):
    slide=ppt.slides[0]
    #康奈尔笔记
    # left, top, width, height = Cm(4.23), Cm(11.13), Cm(3.89), Cm(0.74)
    textBox = slide.shapes.add_textbox(left=left, top=top, width=width, height=height)

    textBox.fill.background()  #文本框填充透明
    # 文本框边框样式调整
    textBox.line.fill.background()   #边框透明
    # line = textBox.line
    # line.color.rgb = RGBColor(0, 255, 0)
    # line.width = Cm(0.1)

    # 获取文本框对象
    tf = textBox.text_frame

    # 文本框样式调整
    tf.margin_bottom = Cm(0.1)  # 下边距
    tf.margin_left = 0  # 左边距
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP  # 对齐文本方式：顶端对齐
    tf.word_wrap = True  # 文本框的文字自动对齐

    # 设置内容
    tf.paragraphs[0].text = text

    # 字体样式调整
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT  # 对齐方式：左对齐
    tf.paragraphs[0].font.name = '微软雅黑'  # 字体名称
    tf.paragraphs[0].font.bold = False  # 是否加粗
    tf.paragraphs[0].font.italic = False  # 是否斜体
    tf.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # 字体颜色
    tf.paragraphs[0].font.size = Pt(12)  # 字体大小
    textBox.click_action.target_slide=ppt.slides[page]

first_Page_hypro(Cm(4.23), Cm(11.13), Cm(3.89), Cm(0.74),"1、康奈尔笔记页",康奈尔笔记开始页)
first_Page_hypro(Cm(4.23), Cm(11.83), Cm(3.89), Cm(0.74),"2、股票买卖记录页",股票开始页)
first_Page_hypro(Cm(4.23), Cm(12.53), Cm(5.62), Cm(0.74),"3、野营、旅游用品清单",3)
first_Page_hypro(Cm(4.23), Cm(13.13), Cm(3.89), Cm(0.74),"4、运动计划表",4)
first_Page_hypro(Cm(4.23), Cm(13.73), Cm(3.89), Cm(0.74),"5、摘录页",摘录开始页)
first_Page_hypro(Cm(4.23), Cm(14.33), Cm(3.89), Cm(0.74),"6、横线纸页",横线开始页)
first_Page_hypro(Cm(4.23), Cm(14.93), Cm(3.89), Cm(0.74),"7、米字格纸页",米字格开始页)


# 在所有页面增加统一的内部跳转按键
# 全部页面都需要用的几个页面跳转
def add_Unit_hypro():  # 全部页面都需要用的几个页面跳转
    for slide in range(0, len(ppt.slides)):  # 遍历所有slide页

        slide = ppt.slides[slide]
        #封面首页链接
        # 添加矩形,封面链接
        # 设置位置以及大小
        # 添加形状
        # 设置背景填充
        # 设置背景颜色
        # 设置边框颜色
        # 内部页面跳转
        left, top, width, height = Cm(4.57), Cm(1.41), Cm(1.5), Cm(0.8)
        矩形_封面 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        矩形_封面.fill.background()  # 填充透明
        # 矩形_封面.fill.solid()
        # 矩形_封面.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1    #形状设置为透明
        # 矩形_封面.fill.fore_color.brightness = -0.25
        # 矩形_封面.line.color.rgb = RGBColor(34, 134, 165)  # 边框颜色
        # 矩形_封面.line.color.brightness = 0  # 边框亮度 在-1和1之间浮点,-1黑色,1白色
        # 矩形_封面.line.dash_style = None  #边框线型
        矩形_封面.line.fill.background()  # 边框透明
        矩形_封面.rotation = -20  # 逆时针旋转30度
        矩形_封面.click_action.target_slide = ppt.slides[0]
        矩形_封面.name = '封面'  # 此shape名字为封面
        # 年历链接

        left, top, width, height = Cm(6.62), Cm(1.31), Cm(1.5), Cm(0.8)

        矩形_年历 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)

        矩形_年历.fill.background()
        矩形_年历.line.fill.background()  # 边框透明
        矩形_年历.click_action.target_slide = ppt.slides[1]
        # 年计划链接

        left, top, width, height = Cm(8.79), Cm(1.43), Cm(1.5), Cm(0.8)

        矩形_年计划 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)

        矩形_年计划.fill.background()
        矩形_年计划.line.fill.background()  # 边框透明
        矩形_年计划.click_action.target_slide = ppt.slides[2]
        # 附录链接

        left, top, width, height = Cm(16.4), Cm(0.98), Cm(1.16), Cm(1.23)

        矩形_附录 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)

        矩形_附录.fill.background()
        矩形_附录.line.fill.background()  # 边框透明
        矩形_附录.click_action.target_slide = ppt.slides[3]  # 附录页
        # 方格纸链接

        left, top, width, height = Cm(18.28), Cm(1.65), Cm(0.77), Cm(0.56)

        矩形_方格纸 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)

        矩形_方格纸.fill.background()
        矩形_方格纸.line.fill.background()  # 边框透明
        矩形_方格纸.click_action.target_slide = ppt.slides[方格纸开始页]
        # 5月超链接
        left, top, width, height = Cm(0.63), Cm(13.93), Cm(0.74), Cm(1.96)

        矩形_5月 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)

        矩形_5月.fill.background()
        矩形_5月.line.fill.background()  # 边框透明
        矩形_5月.click_action.target_slide = ppt.slides[月计划起始页]
        矩形_5月_b = ppt.slides[月计划起始页].shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        矩形_5月_b.fill.solid()
        矩形_5月_b.fill.fore_color.rgb = RGBColor(184,134,11)
        矩形_5月_b.fill.fore_color.brightness = 0.3
        矩形_5月_b.line.fill.background()
        # 6月超链接
        left, top, width, height = Cm(0.59), Cm(16.48), Cm(0.74), Cm(1.96)

        矩形_6月 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)

        矩形_6月.fill.background()
        矩形_6月.line.fill.background()  # 边框透明
        矩形_6月.click_action.target_slide = ppt.slides[月计划起始页+1]
        # 7月超链接
        left, top, width, height = Cm(28.32), Cm(3.8), Cm(0.74), Cm(1.96)

        矩形_7月 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)

        矩形_7月.fill.background()
        矩形_7月.line.fill.background()  # 边框透明
        矩形_7月.click_action.target_slide = ppt.slides[月计划起始页+2]
        # 8月超链接
        left, top, width, height = Cm(28.32), Cm(6.29), Cm(0.74), Cm(1.96)

        矩形_8月 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)

        矩形_8月.fill.background()
        矩形_8月.line.fill.background()  # 边框透明
        矩形_8月.click_action.target_slide = ppt.slides[月计划起始页+3]

        # 9月超链接
        left, top, width, height = Cm(28.32), Cm(8.79), Cm(0.74), Cm(1.96)

        矩形_9月 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)

        矩形_9月.fill.background()
        矩形_9月.line.fill.background()  # 边框透明
        矩形_9月.click_action.target_slide = ppt.slides[月计划起始页+4]
        # 10月超链接
        left, top, width, height = Cm(28.32), Cm(11.29), Cm(0.74), Cm(1.96)

        矩形_10月 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)

        矩形_10月.fill.background()
        矩形_10月.line.fill.background()  # 边框透明
        矩形_10月.click_action.target_slide = ppt.slides[月计划起始页+5]
        # 11月超链接
        left, top, width, height = Cm(28.32), Cm(13.81), Cm(0.74), Cm(1.96)

        矩形_11月 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)

        矩形_11月.fill.background()
        矩形_11月.line.fill.background()  # 边框透明
        矩形_11月.click_action.target_slide = ppt.slides[月计划起始页+6]
        # 12月超链接
        left, top, width, height = Cm(28.32), Cm(16.48), Cm(0.74), Cm(1.96)

        矩形_12月 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)

        矩形_12月.fill.background()
        矩形_12月.line.fill.background()  # 边框透明
        矩形_12月.click_action.target_slide = ppt.slides[月计划起始页+7]


#

# '''月计划页面增加到日计划的内部页面跳转,day_num是当月起始页面号'''
def Add_Monthtarget_hyper(month, day_num):  # day_num,每个月的日计划页面起始页

    for dd in range(11, 53):  # 月计划页占位符索引号11到53
        if slide.placeholders[dd].text != '':
            shape = slide.placeholders[dd]
            # nums=51+month_days(2022,month)
            shape.click_action.target_slide = ppt.slides[day_num]
            print(f'月计划页中日历占位符{dd}跳转链接到{day_num}页')
            day_num += 1


# 在月计划页设置5页菜谱和1页习惯打卡内部跳转
def add_menuitem_hyper(month):
    numM = 0
    slide.placeholders[53].text = '壹'
    slide.placeholders[54].text = '贰'
    slide.placeholders[55].text = '叁'
    slide.placeholders[56].text = '肆'
    slide.placeholders[57].text = '伍'
    slide.placeholders[58].text = '习惯打卡'
    for aa in range(53, 58):  # 53-57号文本占位符是对应5个餐谱位置
        shape = slide.placeholders[aa]
        # if month == 5:
        #     shape.click_action.target_slide = ppt.slides[菜谱开始页 + numM + (month - 5) * 5 ]
        #     print(f'{month}月份菜谱占位符{aa}设置内部跳转到第{菜谱开始页 + numM + (month - 5) * 5 }页')
        # else:
        shape.click_action.target_slide = ppt.slides[菜谱开始页 + numM + (month - 5) * 5]
        print(f'{month}月份菜谱占位符{aa}设置内部跳转到第{菜谱开始页 + numM + (month - 5) * 5}页')
        numM += 1
    slide.placeholders[58].click_action.target_slide = ppt.slides[习惯养成开始页 + month - 5 ]  # 习惯打卡位置,58号占位符
    print(f'{month}月份习惯打卡占位符58号设置内部跳转到第{习惯养成开始页 + month - 5 }页,(ppt页数-1)')



# 在第4到12页的月计划页面中,增加月份书签和日历
# # 定义月份字符数组
# List_Month = ['五月', '六月', '七月', '八月', '九月', '十月', '十一月', '十二月']
# # 只指定的页面输入月份及日历
# 月计划起始页 = 3  # 月计划起始页4页
# 月份 = 5

# 操作月计划页面:在第4到12页的月计划页面中,增加月份书签和日历,增加菜单和习惯打卡页链接
for xxx in range(len(List_Month)):  # 遍历月份字符数组
    索引号 = 11  # 占位符起始编号
    print(f'运行第{xxx + 1}次')
    slide = ppt.slides[月计划起始页]
    slide.placeholders[10].text = List_Month[xxx]
    print(f'幻灯片.placeholders-10号为{slide.placeholders[10].text},在第{月计划起始页}页幻灯片(0是第一页)')
    # 获得日历矩阵
    Day_Freamdata = calendar.monthcalendar(year, 月份)
    print(f'日历{Day_Freamdata}')
    # 遍历日历矩阵,将数据写入占位符
    if 月份 == 5 or 月份 == 10:
        行数 = 6
    else:
        行数 = 5
    # 遍历月历
    for 行号 in range(0, 行数):  # 不同月行数不同,5月和10月是六行,6,7,8,9,11,12月是五行
        for 列号 in range(0, 7):  # 7列
            slide.placeholders[索引号].text = str(Day_Freamdata[行号][列号])
            if slide.placeholders[索引号].text == '0':  # 日历内0的字符全部用空格代替
                slide.placeholders[索引号].text = ''
            # print(f'幻灯片.placeholders{索引号}为{slide.placeholders[索引号].text}')
            索引号 += 1
    # 遍历月计划中日历表,在日历表每个日期增加跳转到对应日计划页面
    if 月份 == 12:
        days_plan -= 1
    Add_Monthtarget_hyper(月份, days_plan)
    print(f'{月份}月days_plan现在是{days_plan}')
    月计划起始页 += 1
    月份 += 1
    if 月份 == 7 or 月份 == 10:  # 6,9,11月是30天,记得从下个月1号起
        days_plan += 30
    else:
        days_plan += 31  # 其它月是31天,记得从下个月1号起
    add_menuitem_hyper(xxx + 5)




# 在日计划表内增加文本
def add_daysText_last(a, b, c, mingyan, birth_day):
    # a 是页面，b是倒计时字符串,c是左上角日期字符串,mingyan是名人名言字符串
    # 获取需要添加文字的a页面，b是输入文字
    slide_page = ppt.slides[a]
    slide_page.placeholders[10].text = b  # 索引号为10的文本框占位符,b字符串(倒计时时间)赋值
    slide_page.placeholders[11].text = c  # 索引号为11的文本框占位符,日期星期字符
    slide_page.placeholders[12].text = mingyan  # 索引号为12的文本框占位符,名人名言
    slide_page.placeholders[13].text = birth_day  # 索引号为13的文本框占位符,降临地球天数
    print(f'倒计时时间-日期-名人名言-已来到地球时间已写入第{a}页(0是第一页)')


# """计算到年底的倒计时函数"""
def 倒计时(Start_date):
    # 导入时间模块
    import datetime

    # 输入目标日期，字符串格式
    target_date = '2022-12-31'
    # 输入开始计算时间,字符串格式
    # Start_date = '2022-5-1'
    # 转换为可识别的
    target_date = datetime.datetime.strptime(target_date, '%Y-%m-%d')

    Start_date = datetime.datetime.strptime(Start_date, '%Y-%m-%d')
    # datetime.datetime(2000, 1, 1, 0, 0)

    # 查看目前的时间
    # datetime.datetime.now()
    # datetime.datetime(2022, 4, 5, 14, 15, 27, 748086)
    # 计算间隔时间。这里间隔时间其实最小单位并不是天，而是微秒
    diff = target_date - Start_date
    # 查看时间间隔
    # diff
    return diff.days


def 来地球天数(Start_date):
    # 导入时间模块
    import datetime

    # 输入目标日期，字符串格式
    target_date = '1977-4-28'
    # 输入开始计算时间,字符串格式
    # Start_date = '2022-5-1'
    # 转换为可识别的
    target_date = datetime.datetime.strptime(target_date, '%Y-%m-%d')
    Start_date = datetime.datetime.strptime(Start_date, '%Y-%m-%d')
    # datetime.datetime(2000, 1, 1, 0, 0)
    # 查看目前的时间
    # datetime.datetime.now()
    # datetime.datetime(2022, 4, 5, 14, 15, 27, 748086)
    # 计算间隔时间。这里间隔时间其实最小单位并不是天，而是微秒
    diff = Start_date - target_date

    # 查看时间间隔
    # diff
    return diff.days


# 增加名人名言
with open('new_名人名言.txt', 'r') as files:  # ,encoding='UTF-16'
    mingyan_list = files.readlines()

月计划起始页=5
add_Unit_hypro()

# 日计划页面最下面现实倒计时

# 操作日计划页面:日计划页起始页
print(f'日计划页总数为{len(My_time)}')
print(f'Sour_time为{len(Sour_time)}')
days_plan = 13
for 天 in range(0, len(My_time)):
    if days_plan < 258:  # 日计划页最后一页页码258
        list1 = Sour_time[天]
        list2 = 倒计时(list1)
        birth_day = 来地球天数(list1)
        list3 = My_time[天]
        b = random.randint(0, 1000)  # 随机读取new_名人名言.txt,共1030行
        list4 = mingyan_list[b]
        add_daysText_last(days_plan, str(list2), list3, list4, str(birth_day))
        days_plan = days_plan + 1
        # 天 += 1

# 保存ppt
ppt.save('已生成笔记本.pptx')
