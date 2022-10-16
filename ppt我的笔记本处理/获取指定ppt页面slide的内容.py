import pptx
from pptx import Presentation
prs=Presentation('已生成笔记本.pptx')
编号=0
for i,slide in enumerate(prs.slides):
    # if i==5 or  i==13:    #获取第5页:月计划页,13页:日计划页
    if  i == 5:  # 获取第5页:月计划页,13页:日计划页
        print(f'显示PPT第{i+1}页')
        for shape in  slide.shapes:
            if shape.has_text_frame:
                text_frame=shape.text_frame
                print(f'shapes编号:{编号}:{text_frame.text}    ')
                text="编号"+str(编号)
                slide.shapes[编号].text_frame.text = text
            编号 +=1

    # slide.shapes[2].text_frame.text = '人生苦短111'
prs.save('修改文字.pptx')