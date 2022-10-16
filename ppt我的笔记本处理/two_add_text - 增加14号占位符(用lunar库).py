# """
# 这个文件可以处理“我的笔记本.ptt' 文件以下功能
#     1、月计划页：复制11页，然后在每一页特定位置添加月份文本框，添加当月日历
#     2、每日页:根据从当天到2022年12月31日复制日计划页，并在每页特定位置标记？？？年？月？日星期几。
#     3、在每日页特定位置，打印名人名言。单独从名人名言库中读取。
#     4、月日历链接到对应日计划页
# """

import calendar
import datetime
import random
import sys

import cnlunar
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt, Cm

"""
# a = cnlunar.Lunar() # 为空为当前时间

a = cnlunar.Lunar(datetime.datetime(2019, 2, 4, 22, 30))
dic = {
    '日期': a.date,
    '农历数字': (a.lunarYear, a.lunarMonth, a.lunarDay, '闰' if a.isLunarLeapMonth else ''),
    '农历': '%s %s[%s]年 %s%s' % (a.lunarYearCn, a.year8Char, a.chineseYearZodiac, a.lunarMonthCn, a.lunarDayCn),
    '星期': a.weekDayCn,
    # 未增加除夕
    '今日节日': (a.get_legalHolidays(), a.get_otherHolidays(), a.get_otherLunarHolidays()),
    '八字': ' '.join([a.year8Char, a.month8Char, a.day8Char, a.twohour8Char]),
    '今日节气': a.todaySolarTerms,
    '下一节气': (a.nextSolarTerm, a.nextSolarTermDate, a.nextSolarTermYear),
    '今年节气表': a.thisYearSolarTermsDic,
    '季节': a.lunarSeason,
    '今日时辰': a.twohour8CharList,
    '时辰凶吉': a.get_twohourLuckyList(),
    '生肖冲煞': a.chineseZodiacClash,
    '星座': a.starZodiac,
    '星次': a.todayEastZodiac,
    '彭祖百忌': a.get_pengTaboo(),
    '彭祖百忌精简': a.get_pengTaboo(long=4, delimit='<br>'),
    '十二神': a.get_today12DayOfficer(),
    '廿八宿': a.get_the28Stars(),
    '今日三合': a.zodiacMark3List,
    '今日六合': a.zodiacMark6,
    '今日五行': a.get_today5Elements(),
    '纳音': a.get_nayin(),
    '九宫飞星': a.get_the9FlyStar(),
    '吉神方位': a.get_luckyGodsDirection(),
    '今日胎神': a.get_fetalGod(),
    '神煞宜忌': a.angelDemon,
    '今日吉神': a.goodGodName,
    '今日凶煞': a.badGodName,
    '宜': a.goodThing,
    '忌': a.badThing,
    '时辰经络': a.meridians
}
"""

# 打开已存在ppt
ppt = Presentation('2笔记本_增加页面.pptx')
print(f'这个ppt共有{len(ppt.slides)}页,0为第一页')
'''
每个月至少有28天，所以初始值设成这个月的28号，然后尝试对这个日期加1天，加2天，加3天，什么时候日期跳到下个月了就停止，这样就知道这个月的天数了。
判断一个月有多少天
'''

# 全局变量和常数定义
week_list = ["星期一", "星期二", "星期三", "星期四", "星期五", "星期六", "星期日"]
YEAR = 2022
num = 0
My_time = []  # 生成[2022年5月1日星期日]格式的My_time列表;
Sour_time = []  # 生成[2022-5-1]格式的My_time列表;
日计划起始页 = 13  # 日计划起始页,第14页,还有一个位置有同名参数,记得同时改
List_Month = ['五月', '六月', '七月', '八月', '九月', '十月', '十一月', '十二月']  # 定义月份字符数组
月计划起始页 = 5  # 月计划起始页6页,还有一个位置有同名参数,记得同时改
月份 = 5  # 从5月份开始计算
习惯养成开始页 = 258
菜谱开始页 = 266
股票开始页 = 306
康奈尔笔记开始页 = 336
方格纸开始页 = 366
米字格开始页 = 371
横线开始页 = 376
摘录开始页 = 381
所在行 = sys._getframe().f_lineno


# 判断一个月有多少天
def month_days(year, month):
    begin = datetime.date(year, month, 28)
    for i in range(1, 5):
        end = begin + datetime.timedelta(days=i)
        if end.month > month:
            break
    result = 28 + i - 1
    return result  # 返回一个月的天数int


# 将年月日星期数据存入My_time列表
for i in range(5, 13):  # 5-12月
    for n in range(1, month_days(YEAR, i) + 1):  # 每个月天数
        week = week_list[datetime.date(YEAR, i, n).weekday()]
        My_time.append(str(YEAR) + '年' + str(i) + '月' + str(n) + '日' + week)
        Sour_time.append(str(YEAR) + '-' + str(i) + '-' + str(n))


# 在第一页增加股票,摘抄等链接

def first_Page_hypro(left, top, width, height, text, page):
    slide = ppt.slides[0]
    # 康奈尔笔记
    # left, top, width, height = Cm(4.23), Cm(11.13), Cm(3.89), Cm(0.74)
    textBox = slide.shapes.add_textbox(left=left, top=top, width=width, height=height)
    textBox.fill.background()  # 文本框填充透明
    # 文本框边框样式调整
    textBox.line.fill.background()  # 边框透明
    # line = textBox.line
    # line.color.rgb = RGBColor(0, 255, 0)
    # line.width = Cm(0.1)

    # 获取文本框对象
    tf = textBox.text_frame

    # 文本框样式调整
    tf.margin_bottom = Cm(0.1)  # 下边距
    tf.margin_left = 0  # 左边距
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP  # 对齐文本方式：顶端对齐
    tf.word_wrap = True  # 文本框的文字自动对齐

    # 设置内容
    tf.paragraphs[0].text = text

    # 字体样式调整
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT  # 对齐方式：左对齐
    tf.paragraphs[0].font.name = '方正小标宋简体'  # 字体名称
    tf.paragraphs[0].font.bold = True  # 是否加粗
    tf.paragraphs[0].font.italic = False  # 是否斜体
    tf.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # 字体颜色
    tf.paragraphs[0].font.size = Pt(12)  # 字体大小
    textBox.click_action.target_slide = ppt.slides[page]


# 全部页面都需要用的几个页面跳转(封面,年历,年计划# 附录# 方格纸)等
def all_Unit_hypro(left, top, width, height, name, rotation, slide, hypro):
    left, top, width, height = Cm(left), Cm(top), Cm(width), Cm(height)
    name = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    name.fill.background()  # 填充透明
    # name.fill.solid()
    # name.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1    #形状设置为透明
    # name.fill.fore_color.brightness = -0.25
    # name.line.color.rgb = RGBColor(34, 134, 165)  # 边框颜色
    # name.line.color.brightness = 0  # 边框亮度 在-1和1之间浮点,-1黑色,1白色
    # name.line.dash_style = None  #边框线型
    name.line.fill.background()  # 边框透明
    name.rotation = -20  # 逆时针旋转30度
    name.click_action.target_slide = ppt.slides[hypro]
    name.name = str(name)  # 此shape名字为封面
    # print(f'all_Unit_hypro函数中slide={slide},写入{name.name}')


# '''月计划页面增加到日计划的内部页面跳转,day_num是当月起始页面号'''
def Add_Monthtarget_hyper(month, day_num, slide):  # day_num,每个月的日计划页面起始页
    for dd in range(11, 53):  # 月计划页占位符索引号11到53
        if slide.placeholders[dd].text != '':
            shape = slide.placeholders[dd]
            # nums=51+month_days(YEAR,month)
            shape.click_action.target_slide = ppt.slides[day_num]
            # print(f'月计划页中日历占位符{dd}跳转链接到{day_num}页')
            day_num += 1
    # print(f'Add_monthtarget_hyper()函数中slide={slide}')


# def add_Unit_hypro():  # 全部页面都需要用的几个页面跳转
#     for slide in range(0, len(ppt.slides)):  # 遍历所有slide页
#
#         slide = ppt.slides[slide]
#         #封面首页链接
#         # 添加矩形,封面链接
#         # 设置位置以及大小
#         # 添加形状
#         # 设置背景填充
#         # 设置背景颜色
#         # 设置边框颜色
#         # 内部页面跳转
#         left, top, width, height = Cm(4.57), Cm(1.41), Cm(1.5), Cm(0.8)
#         矩形_封面 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
#         矩形_封面.fill.background()  # 填充透明
#         # 矩形_封面.fill.solid()
#         # 矩形_封面.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1    #形状设置为透明
#         # 矩形_封面.fill.fore_color.brightness = -0.25
#         # 矩形_封面.line.color.rgb = RGBColor(34, 134, 165)  # 边框颜色
#         # 矩形_封面.line.color.brightness = 0  # 边框亮度 在-1和1之间浮点,-1黑色,1白色
#         # 矩形_封面.line.dash_style = None  #边框线型
#         矩形_封面.line.fill.background()  # 边框透明
#         矩形_封面.rotation = -20  # 逆时针旋转30度
#         矩形_封面.click_action.target_slide = ppt.slides[0]
#         矩形_封面.name = '封面'  # 此shape名字为封面
#         # 年历链接
#
#         left, top, width, height = Cm(6.62), Cm(1.31), Cm(1.5), Cm(0.8)
#
#         矩形_年历 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
#
#         矩形_年历.fill.background()
#         矩形_年历.line.fill.background()  # 边框透明
#         矩形_年历.click_action.target_slide = ppt.slides[1]
#         # 年计划链接
#
#         left, top, width, height = Cm(8.79), Cm(1.43), Cm(1.5), Cm(0.8)
#
#         矩形_年计划 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
#
#         矩形_年计划.fill.background()
#         矩形_年计划.line.fill.background()  # 边框透明
#         矩形_年计划.click_action.target_slide = ppt.slides[2]
#         # 附录链接
#
#         left, top, width, height = Cm(16.4), Cm(0.98), Cm(1.16), Cm(1.23)
#
#         矩形_附录 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
#
#         矩形_附录.fill.background()
#         矩形_附录.line.fill.background()  # 边框透明
#         矩形_附录.click_action.target_slide = ppt.slides[3]  # 附录页
#         # 方格纸链接
#
#         left, top, width, height = Cm(18.28), Cm(1.65), Cm(0.77), Cm(0.56)
#
#         矩形_方格纸 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
#
#         矩形_方格纸.fill.background()
#         矩形_方格纸.line.fill.background()  # 边框透明
#         矩形_方格纸.click_action.target_slide = ppt.slides[方格纸开始页]
#         # 5月超链接
#         left, top, width, height = Cm(0.63), Cm(13.93), Cm(0.74), Cm(1.96)
#
#         # 矩形_5月 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
#         矩形_5月 = slide.shapes.add_textbox( left, top, width, height)
#
#         # 调整文本框背景颜色
#         textBoxFill = 矩形_5月.fill
#         textBoxFill.background()
#         # textBoxFill.fore_color.rgb = RGBColor(187, 255, 255)
#         # 文本框边框样式调整
#         line = 矩形_5月.line
#         line.fill.background()
#         # line.color.rgb = RGBColor(0, 255, 0)
#         # line.width = Cm(0.1)
#         # 获取文本框对象
#         tf = 矩形_5月.text_frame
#         # 文本框样式调整
#         tf.margin_bottom = Cm(0.1)  # 下边距
#         tf.margin_left = 0  # 左边距
#         tf.vertical_anchor = MSO_VERTICAL_ANCHOR.BOTTOM  # 对齐文本方式：底端对齐
#         tf.word_wrap = True  # 文本框的文字自动对齐
#         tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
#         # 设置内容
#         tf.paragraphs[0].text = '五月'
#
#         # 字体样式调整
#         tf.paragraphs[0].alignment = PP_ALIGN.CENTER  # 对齐方式  DISTRIBUTE JUSTIFY JUSTIFY_LOW  LEFT RIGHT THAI_DISTRIBUTE MIXED
#         tf.paragraphs[0].font.name = '华文隶书'  # 字体名称
#         tf.paragraphs[0].font.bold = True  # 是否加粗
#         tf.paragraphs[0].font.italic = False  # 是否斜体
#         tf.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)  # 字体颜色
#         tf.paragraphs[0].font.size = Pt(18)  # 字体大小
#         # 矩形_5月.text="五月"
#         # 矩形_5月.font.size=Pt(15)
#         # 矩形_5月.fill.background()
#         # 矩形_5月.line.fill.background()  # 边框透明
#         矩形_5月.click_action.target_slide = ppt.slides[月计划起始页]
#         # 6月超链接
#         left, top, width, height = Cm(0.59), Cm(16.48), Cm(0.74), Cm(1.96)
#
#         矩形_6月 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
#
#         矩形_6月.fill.background()
#         矩形_6月.line.fill.background()  # 边框透明
#         矩形_6月.click_action.target_slide = ppt.slides[月计划起始页+1]
#         # 7月超链接
#         left, top, width, height = Cm(28.32), Cm(3.8), Cm(0.74), Cm(1.96)
#
#         矩形_7月 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
#
#         矩形_7月.fill.background()
#         矩形_7月.line.fill.background()  # 边框透明
#         矩形_7月.click_action.target_slide = ppt.slides[月计划起始页+2]
#         # 8月超链接
#         left, top, width, height = Cm(28.32), Cm(6.29), Cm(0.74), Cm(1.96)
#
#         矩形_8月 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
#
#         矩形_8月.fill.background()
#         矩形_8月.line.fill.background()  # 边框透明
#         矩形_8月.click_action.target_slide = ppt.slides[月计划起始页+3]
#
#         # 9月超链接
#         left, top, width, height = Cm(28.32), Cm(8.79), Cm(0.74), Cm(1.96)
#
#         矩形_9月 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
#
#         矩形_9月.fill.background()
#         矩形_9月.line.fill.background()  # 边框透明
#         矩形_9月.click_action.target_slide = ppt.slides[月计划起始页+4]
#         # 10月超链接
#         left, top, width, height = Cm(28.32), Cm(11.29), Cm(0.74), Cm(1.96)
#
#         矩形_10月 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
#
#         矩形_10月.fill.background()
#         矩形_10月.line.fill.background()  # 边框透明
#         矩形_10月.click_action.target_slide = ppt.slides[月计划起始页+5]
#         # 11月超链接
#         left, top, width, height = Cm(28.32), Cm(13.81), Cm(0.74), Cm(1.96)
#
#         矩形_11月 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
#
#         矩形_11月.fill.background()
#         矩形_11月.line.fill.background()  # 边框透明
#         矩形_11月.click_action.target_slide = ppt.slides[月计划起始页+6]
#         # 12月超链接
#         left, top, width, height = Cm(28.32), Cm(16.48), Cm(0.74), Cm(1.96)
#
#         矩形_12月 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
#
#         矩形_12月.fill.background()
#         矩形_12月.line.fill.background()  # 边框透明
#         矩形_12月.click_action.target_slide = ppt.slides[月计划起始页+7]
# 每页对应的1-12月跳转链接
# 添加月份跳转
def month_hypro(left, top, width, height, r, b, g, names, hypro, slide):
    # 添加月份跳转
    # 左,顶,宽,高位置,r,b,g颜色,文本框内容,跳转页面
    # 月份超链接
    titles = names
    # print(f'338行---------------name={names}---slide={slide}')
    left, top, width, height = Cm(left), Cm(top), Cm(width), Cm(height)
    names = slide.shapes.add_textbox(left, top, width, height)
    # 调整文本框背景颜色
    textBoxFill = names.fill
    textBoxFill.background()
    names.line.fill.background()

    # 获取文本框对象
    tf = names.text_frame
    # 文本框样式调整
    tf.margin_bottom = Cm(0.1)  # 下边距
    tf.margin_left = 0  # 左边距
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.BOTTOM  # 对齐文本方式：底端对齐
    tf.word_wrap = True  # 文本框的文字自动对齐
    tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT  # SHAPE_TO_FIT_TEXT调整形状高度和可能的宽度以适合文本。请注意，此设置与 TextFrame.word_wrap 属性设置交互。如果打开自动换行，只会调整形状的高度；软换行符将用于水平适应文本。
    # tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE  #根据需要减小字体大小以适应形状内的文本。
    # 设置内容
    tf.paragraphs[0].text = titles

    # 字体样式调整
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT  # 对齐方式  DISTRIBUTE JUSTIFY JUSTIFY_LOW  LEFT RIGHT THAI_DISTRIBUTE MIXED
    tf.paragraphs[0].font.name = '方正小标宋繁体'  # 字体名称
    tf.paragraphs[0].font.bold = True  # 是否加粗
    tf.paragraphs[0].font.italic = False  # 是否斜体
    tf.paragraphs[0].font.color.rgb = RGBColor(r, g, b)  # 字体颜色
    # print(f'364行:r,g,b={r}-{g}-{b}')
    tf.paragraphs[0].font.size = Pt(19)  # 字体大小
    names.click_action.target_slide = ppt.slides[hypro]
    # print(f'366行------------写入{titles}-------跳转到{hypro}')


# 在月计划页设置5页菜谱和1页习惯打卡内部跳转
def add_menuitem_hyper(month, slide):
    numM = 0
    slide.placeholders[53].text = '壹'
    slide.placeholders[54].text = '贰'
    slide.placeholders[55].text = '叁'
    slide.placeholders[56].text = '肆'
    slide.placeholders[57].text = '伍'
    slide.placeholders[58].text = '习惯打卡'
    for aa in range(53, 58):  # 53-57号文本占位符是对应5个餐谱位置
        shape = slide.placeholders[aa]
        shape.click_action.target_slide = ppt.slides[菜谱开始页 + numM + (month - 5) * 5]  # 从5月起
        # print(f'{month}月份菜谱占位符{aa}设置内部跳转到第{菜谱开始页 + numM + (month - 5) * 5}页')
        numM += 1
    slide.placeholders[58].click_action.target_slide = ppt.slides[习惯养成开始页 + month - 5]  # 习惯打卡位置,58号占位符
    # print(f'{month}月份习惯打卡占位符58号设置内部跳转到第{习惯养成开始页 + month - 5}页,(ppt页数-1)')


# 操作月计划页面:在第4到12页的月计划页面中,增加月份书签和日历,增加菜单和习惯打卡页链接
def month_page_add():
    global 月计划起始页, 月份, 日计划起始页
    for xxx in range(len(List_Month)):  # 遍历月份字符数组
        索引号 = 11  # 占位符起始编号
        # print(f'运行第{xxx + 1}次')
        slide = ppt.slides[月计划起始页]
        slide.placeholders[10].text = List_Month[xxx]
        # print(f'幻灯片.placeholders-10号为{slide.placeholders[10].text},在第{月计划起始页}页幻灯片(0是第一页)')
        # 获得日历矩阵
        Day_Freamdata = calendar.monthcalendar(YEAR, 月份)
        # print(f'日历{Day_Freamdata}')
        # 遍历日历矩阵,将数据写入占位符
        if 月份 == 5 or 月份 == 10:
            行数 = 6
        else:
            行数 = 5
        # 遍历月历
        for 行号 in range(0, 行数):  # 不同月行数不同,5月和10月是六行,6,7,8,9,11,12月是五行
            for 列号 in range(0, 7):  # 7列
                slide.placeholders[索引号].text = str(Day_Freamdata[行号][列号])
                if slide.placeholders[索引号].text == '0':  # 日历内0的字符全部用空格代替
                    slide.placeholders[索引号].text = ''
                # print(f'幻灯片.placeholders{索引号}为{slide.placeholders[索引号].text}')
                索引号 += 1
        # 遍历月计划中日历表,在日历表每个日期增加跳转到对应日计划页面
        if 月份 == 12:
            日计划起始页 -= 1
        Add_Monthtarget_hyper(月份, 日计划起始页, slide)
        print(f'{月份}月日计划起始页现在是{日计划起始页}')
        月计划起始页 += 1
        月份 += 1
        if 月份 == 7 or 月份 == 10:  # 6,9,11月是30天,记得从下个月1号起
            日计划起始页 += 30
        else:
            日计划起始页 += 31  # 其它月是31天,记得从下个月1号起
        add_menuitem_hyper(xxx + 5, slide)


# 在日计划表内增加文本
def add_daysText_last(a, b, c, mingyan, birth_day, text, yinli):
    # a 是页面，b是倒计时字符串,c是左上角日期字符串,mingyan是名人名言字符串
    # 获取需要添加文字的a页面，b是输入文字
    slide_page = ppt.slides[a]
    slide_page.placeholders[10].text = b  # 索引号为10的文本框占位符,b字符串(倒计时时间)赋值
    slide_page.placeholders[11].text = c  # 索引号为11的文本框占位符,日期星期字符
    slide_page.placeholders[12].text = mingyan  # 索引号为12的文本框占位符,名人名言
    slide_page.placeholders[13].text = birth_day  # 索引号为13的文本框占位符,降临地球天数
    slide_page.placeholders[14].text = text  # 索引号为14的文本框占位符,段落,左下角可以放笑话,每日一句英语或者其它文字多一点的内容
    slide_page.placeholders[15].text = yinli  # 显示农历或者阴历日期
    # print(f'倒计时时间-日期-名人名言-已来到地球时间已写入第{a}页(0是第一页)')


# 在习惯养成页添加年/月信息
def add_xiguan_date(page, year, month):  # 在习惯养成页添加年/月信息
    slide_page = ppt.slides[page]
    slide_page.placeholders[10].text = str(year) + '年'  # 索引号为10的文本框占位符,显示年
    slide_page.placeholders[11].text = str(month)  # 索引号为11的文本框占位符,显示月


# 在每周菜谱页添加月及第几周信息
def add_item_date(page, month, week):  # 在每周菜谱页添加月及第几周信息
    slide_page = ppt.slides[page]
    slide_page.placeholders[10].text = str(month)  # 索引号为10的文本框占位符,显示月
    if week == 1: week = "一"
    if week == 2: week = "二"
    if week == 3: week = "三"
    if week == 4: week = "四"
    if week == 5: week = "五"
    slide_page.placeholders[11].text = '第' + week + "周"  # 索引号为11的文本框占位符,显示第几周


# """计算到年底的倒计时函数"""
def 倒计时(Start_date):
    # 导入时间模块
    import datetime

    # 输入目标日期，字符串格式
    target_date = '2022-12-31'
    # 输入开始计算时间,字符串格式
    # Start_date = '2022-5-1'
    # 转换为可识别的
    target_date = datetime.datetime.strptime(target_date, '%Y-%m-%d')

    Start_date = datetime.datetime.strptime(Start_date, '%Y-%m-%d')
    # datetime.datetime(2000, 1, 1, 0, 0)

    # 查看目前的时间
    # datetime.datetime.now()
    # datetime.datetime(2022, 4, 5, 14, 15, 27, 748086)
    # 计算间隔时间。这里间隔时间其实最小单位并不是天，而是微秒
    diff = target_date - Start_date
    # 查看时间间隔
    # diff
    return diff.days


def 来地球天数(Start_date):
    # 导入时间模块
    import datetime

    # 输入目标日期，字符串格式
    target_date = '1977-4-28'
    # 输入开始计算时间,字符串格式
    # Start_date = '2022-5-1'
    # 转换为可识别的
    target_date = datetime.datetime.strptime(target_date, '%Y-%m-%d')
    Start_date = datetime.datetime.strptime(Start_date, '%Y-%m-%d')
    # datetime.datetime(2000, 1, 1, 0, 0)
    # 查看目前的时间
    # datetime.datetime.now()
    # datetime.datetime(2022, 4, 5, 14, 15, 27, 748086)
    # 计算间隔时间。这里间隔时间其实最小单位并不是天，而是微秒
    diff = Start_date - target_date

    # 查看时间间隔
    # diff
    return diff.days


# 月计划起始页 = 5


# 写入每页都包含的链接内容
def write_all_page():
    additon_list = []
    with open('addition.csv', 'r', encoding='UTF8') as files:
        additon_list = files.readlines()
    for slide in range(0, len(ppt.slides)):  # 遍历所有slide页
        slide = ppt.slides[slide]
        # print(f'slide={slide}')
        # 月计划起始页 = 5
        # 几个页面跳转(封面,年历,年计划# 附录# 方格纸)等
        for index in range(13, 18):  # 第13行到18行
            css_table_1 = additon_list[index].split(',')
            # print(css_table_1)
            left, top, width, height = float(css_table_1[1]), float(css_table_1[2]), float(css_table_1[3]), float(
                css_table_1[4])
            name = str(css_table_1[0])
            hypro = int(css_table_1[11].strip())
            # print(f'第505行代码---------------name={name}')
            all_Unit_hypro(left, top, width, height, name, int(css_table_1[5]), slide=slide,
                           hypro=hypro)  # 用pandas取csv文件内容

        # 每页对应的1-12月跳转链接
        for index1 in range(5, 13):
            css_table_2 = additon_list[index1].split(',')
            # print(css_table_2)
            left, top, width, height = float(css_table_2[1]), float(css_table_2[2]), float(css_table_2[3]), float(
                css_table_2[4])
            r, g, b = css_table_2[5], css_table_2[6], css_table_2[7]
            # print(f'r={r}类型为{type(r)}')    #字符串类型,删除头尾空格,转换为整数型
            r, g, b = int(r.strip()), int(g.strip()), int(b.strip())
            # print(f'r={r}类型为{type(r)}')
            name = str(css_table_2[0])
            # print(f'第513行代码---------------name={name}')
            month_hypro(left, top, width, height, r, g, b, name, int(css_table_2[11].strip()), slide=slide)
            # 月计划起始页 += 1
        # 在月计划页变化月标签字体斜体,加粗,变色
    # 在月计划页改变对应月份链接文字的颜色
    # global 月计划起始页
    # for month in range(月份, 13):
    #     css_table_3 = additon_list[month].split(',')
    #     slide = ppt.slides[月计划起始页]
    #     left, top, width, height = float(css_table_3[1]), float(css_table_3[2]), float(css_table_3[3]), float(
    #         css_table_3[4])
    #     r, g, b = css_table_3[8], css_table_3[9], css_table_3[10]
    #     # print(f'r={r}类型为{type(r)}')    #字符串类型,删除头尾空格,转换为整数型
    #     r, g, b = int(r.strip()), int(g.strip()), int(b.strip())
    #     # print(f'r={r}类型为{type(r)}')
    #     name = str(css_table_3[0])
    #     # print(f'第513行代码---------------name={name}')
    #     month_hypro(left, top, width, height, r, g, b, name, int(css_table_3[11].strip()), slide=slide)
    #     月计划起始页 += 1


# 在日计划页写入占位符信息
def write_day_page():
    # 增加名人名言
    with open('new_名人名言.txt', 'r') as files:  # ,encoding='UTF8'
        mingyan_list = files.readlines()
    with open('左下角文本.txt', 'r', encoding='utf8') as f:  # ,encoding='UTF8'
        左下角文本 = f.readlines()
    # 操作日计划页面:日计划页起始页
    日计划起始页 = 13
    for 天 in range(0, len(My_time)):
        if 日计划起始页 < 258:  # 日计划页最后一页页码258
            list1 = Sour_time[天]
            list2 = 倒计时(list1)
            birth_day = 来地球天数(list1)
            list3 = My_time[天]
            b = random.randint(0, 1000)  # 随机读取new_名人名言.txt,共1030行
            list4 = mingyan_list[b]
            # print(datetime.datetime.strptime(list1, '%Y-%m-%d'))
            yinli = list1.split('-')
            # a = cnlunar.lunar(datetime.datetime.strptime(list1, '%Y-%m-%d'))
            a = cnlunar.Lunar(datetime.datetime(int(yinli[0]), int(yinli[1]), int(yinli[2]), 0, 0))
            阴历内容 = '%s %s[%s]年 %s%s' % (
                a.lunarYearCn, a.year8Char, a.chineseYearZodiac, a.lunarMonthCn, a.lunarDayCn)  # 阳历转农历
            # 阴历内容=阴历内容+'  节日:'+a.get_legalHolidays() +' 节气:'+a.todaySolarTerms+"宜:"+''.join(a.goodThing)
            阴历内容 = 阴历内容 + '   ' + a.get_legalHolidays().replace('无', '') + '      ' + a.todaySolarTerms.replace('无', '')
            # print(f'阴历信息:{阴历内容}')
            add_daysText_last(日计划起始页, str(list2), list3, list4, str(birth_day), 左下角文本[日计划起始页], 阴历内容)
            日计划起始页 = 日计划起始页 + 1
            # 天 += 1


# 改变指定页面文字标签的颜色和斜体
def change_text_color(sharp_strat_num, page, num):
    # sharp_start_num: 要修改的sharp第一个编号
    # page:页面编号
    # num:从0为第一个编号

    slide = ppt.slides[page]

    # print(f'{sys._getframe().f_lineno}:修改第{page}页shape编号为{sharp_strat_num+4+num}的字体')
    # slide.shapes[sharp_strat_num+4+num].text_frame.paragraphs[0].text = '人生苦短人生苦短'   #5月
    slide.shapes[sharp_strat_num + 4 + num].text_frame.paragraphs[0].font.italic = True  # 斜体
    slide.shapes[sharp_strat_num + 4 + num].text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # 改变字体颜色
    # print(f'{sys._getframe().f_lineno}:修改完毕')


def main():
    global 习惯养成开始页, YEAR, 月份
    first_Page_hypro(Cm(4.23), Cm(11.13), Cm(3.89), Cm(0.74), "1、康奈尔笔记页", 康奈尔笔记开始页)
    first_Page_hypro(Cm(4.23), Cm(11.83), Cm(3.89), Cm(0.74), "2、股票买卖记录页", 股票开始页)
    first_Page_hypro(Cm(4.23), Cm(12.53), Cm(5.62), Cm(0.74), "3、野营、旅游用品清单", 3)
    first_Page_hypro(Cm(4.23), Cm(13.13), Cm(3.89), Cm(0.74), "4、运动计划表", 4)
    first_Page_hypro(Cm(4.23), Cm(13.73), Cm(3.89), Cm(0.74), "5、摘录页", 摘录开始页)
    first_Page_hypro(Cm(4.23), Cm(14.33), Cm(3.89), Cm(0.74), "6、横线纸页", 横线开始页)
    first_Page_hypro(Cm(4.23), Cm(14.93), Cm(3.89), Cm(0.74), "7、米字格纸页", 米字格开始页)
    month_page_add()  # 操作月计划页面:在第4到12页的月计划页面中,增加月份书签和日历,增加菜单和习惯打卡页链接
    write_all_page()
    write_day_page()
    global 习惯养成开始页, YEAR, 菜谱开始页
    for aa in range(0, 8):  # 5-12页循环8次
        add_xiguan_date(习惯养成开始页, YEAR, List_Month[aa])

        for bb in range(0, 5):  # 每月5页
            add_item_date(菜谱开始页, List_Month[aa], bb + 1)
            菜谱开始页 += 1
        习惯养成开始页 += 1


if __name__ == '__main__':
    # 月计划起始页,日计划起始页=5,13

    num = 0
    print(f'{sys._getframe().f_lineno}:年月日格式的日期数组My_time是{My_time}')
    print(f'{sys._getframe().f_lineno}:日期格式的年月日数组Sour_time是{Sour_time}')
    print(f'{sys._getframe().f_lineno}:日计划页总数为{len(My_time)}')
    print(f'{sys._getframe().f_lineno}:Sour_time数列内数量为{len(Sour_time)}')
    main()
    print(f"{sys._getframe().f_lineno}:begin start change text color")
    日计划起始页 = 13
    月计划起始页 = 5
    # 月计划页面月标签字体改色+斜体
    for page in range(月计划起始页, 13):  # 从月计划开始页到第12页
        # page页,从shape第50编号起,改变num+1个标签
        print(f'{sys._getframe().f_lineno}:月计划起始页={月计划起始页}')
        change_text_color(50 + page - 5, page, 0)  # 5月开始
        # print(f'{sys._getframe().f_lineno}:运行第{num +1}次')
        # num+=1
    # 改变日计划页面中对应月标签颜色+斜体
    all_days = 0
    num = 0
    for nums in range(5, 13):  # 5月起
        月日开始 = 日计划起始页 + all_days
        days = month_days(YEAR, nums)
        all_days = all_days + days
        月日结束 = 日计划起始页 + all_days - 1
        print(f'{sys._getframe().f_lineno}:{nums}月日开始{月日开始}:{nums}月日结束{月日结束}')
        print(f'{sys._getframe().f_lineno}:开始修改{nums}月份的日计划页面对应月份链接,从{月日开始}页到{月日结束}页')
        for page in range(月日开始, 月日结束 + 1):
            print(f'{sys._getframe().f_lineno}:修改第{page}页')
            sharp_start_num = 7 + nums - 5
            change_text_color(sharp_start_num, page, 0)
        # num += 1

# 保存ppt
ppt.save('已生成笔记本.pptx')
#
# #ppt导出为pdf格式
# import office
# office.ppt.ppt2pdf(path='E:\\python workbook\\ppt我的笔记本处理')
