from pptx import Presentation
import calendar
# import pptx.enum
# from pptx.opc.constants import RELATIONSHIP_TYPE as RT
import pandas as pd
文件 = Presentation('3增加文字.pptx')
幻灯片=文件.slides[16] #操作第17页
# 数据=pd.read_excel("日历.xlsx",sheet_name='5月',header=0,dtype="str")
# print(数据)
# 数据.astype('str').dtypes
# 数据.fillna(' ',inplace=True)  #用空字符代替NaN字符,并且重新写入原有数据表
# print(数据)

# 用python自己生成日历表
year=2022
month_data=5
数据=calendar.monthcalendar(year,month_data)  #返回矩阵

print(数据)

print(数据[2][3])




# print(str(数据.iloc[ [ 0 ], [ 0 ] ]))
# print(数据.shape)
def Month_list():
    索引号=10
    for 列号 in range(0,7):
        for 行号 in range(0,6):  #要第2行到第7行数据
        # print(数据.iloc[ 行号 ], [ 列号 ])
        # print(type(数据.iloc[ 行号 ], [ 列号 ]))
        # 幻灯片.placeholders[ 索引号 ].text = str(数据.iat[行号, 列号])
            幻灯片.placeholders[ 索引号 ].text = str(数据[行号][列号])
            索引号 +=1


#
#
# 链接=幻灯片.placeholders[ 15].text_frame.paragraphs[0].add_run()
# 链接.text =''
# rPr = r._r.get_or_add_rPr()
# hlinkClick = rPr.add_hlinkClick(rId)
#
# hlinkClick.set('action', 'ppaction://hlinksldjump')
#
#
#
#




文件.save("4导入日历表.pptx")

#
# """
# 增加超链接
# 先向文本添加"run",添加要显示文本,然后添加链接地址
# cell = table.cell(0,1)
# cell_link = cell.text_frame.paragraphs[0].add_run()
# cell_link.text = "Text to display"
# hlink = cell_link.hyperlink
# hlink.address = "https://www.somwehere.com"
#
# """