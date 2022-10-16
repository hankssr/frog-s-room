"""
这个文件可以处理“我的笔记本.ptt' 文件以下功能
    1、月计划页：复制11页，然后在每一页特定位置添加月份文本框，添加当月日历
    2、每日页:根据从当天到2022年12月31日复制日计划页，并在每页特定位置标记？？？年？月？日星期几。
    3、在每日页特定位置，打印名人名言。单独从名人名言库中读取。
    4、月日历链接到对应日计划页
"""
import datetime
import win32com
from win32com.client import Dispatch
import os
from pptx import Presentation
import time


# """计算到年底的倒计时函数"""
def 倒计时():
    ## 导入时间模块
    # import datetime

    ## 输入目标日期，字符串格式
    Target_date = '2022-12-31'
    # 输入开始计算时间,字符串格式
    Start_date = '2022-5-1'
    ## 转换为可识别的
    Target_date = datetime.datetime.strptime(Target_date, '%Y-%m-%d')
    Target_date
    Start_date = datetime.datetime.strptime(Start_date, '%Y-%m-%d')
    # datetime.datetime(2000, 1, 1, 0, 0)

    ## 查看目前的时间
    # datetime.datetime.now()

    # datetime.datetime(2022, 4, 5, 14, 15, 27, 748086)

    ## 计算间隔时间。这里间隔时间其实最小单位并不是天，而是微秒
    diff = Target_date - Start_date

    ## 查看时间间隔
    diff

    # datetime.timedelta(days=8130, seconds=58676, microseconds=52564)

    ## 提取天数的部分
    return diff.days  ## 其它部分忽略不算

def Copy_ppt(page_num, pageNum):  # 复制指定页面:page_num页号,pageNum复印页数量
    pptSel.Slides(page_num).Copy()
    for i in range(pageNum):
        pptSel.Slides.Paste()

def Del_page():
    ppt = Presentation('笔记本_增加页面temp.pptx')  # 实例化Presentation类，构建一个名为ppt的对象
    sldidss = list(ppt.slides._sldIdLst)  # 创建slideid列表，每个id对应一张ppt

    for index in range(4, 16):  # 删除第5-15张，这个对应的是python列表的操作方法
        ppt.slides._sldIdLst.remove(sldidss[index])
    print(f'已删除aa页,删除后新文件共{len(sldidss)}页')
    ppt.save('2笔记本_增加页面.pptx')  # 保存ppt，可以在jupyter notebook的工作目录下找到

ppt = Dispatch('PowerPoint.Application')
# 或者使用下面的方法，使用启动独立的进程：
# ppt = DispatchEx('PowerPoint.Application')
# 如果不声明以下属性，运行的时候会显示的打开word
ppt.Visible = 1  # 后台运行
ppt.DisplayAlerts = 0  # 不显示，不警告
# 创建新的PowerPoint文档
# pptSel = ppt.Presentations.Add()
# 打开一个已有的PowerPoint文档
pptSel = ppt.Presentations.Open(os.getcwd() + "\\" + "笔记本.pptx")

# 复制模板页,第4页是月计划(编号4)，5-12月，复印7页
# 注意:通过win32com调用复制,第一页索引号是1
Copy_ppt(4, 8)  # 月计划表:8页
pageNums = int(倒计时() + 1)
Copy_ppt(7, pageNums)  # 复制第7页日计划页
time.sleep(8)
Copy_ppt(5, 8)  # 习惯打卡8页
time.sleep(3)
Copy_ppt(6, 40)  # 每周菜单5*8(5-12月)=40页
time.sleep(3)
Copy_ppt(8, 30)  # 康奈尔笔记30页
Copy_ppt(9, 30)  # 股票记录30页
time.sleep(3)
Copy_ppt(12, 5)  # 摘录页:5页
Copy_ppt(13, 5)  # 方格纸:5页
Copy_ppt(14, 5)  # 横线纸:5页
time.sleep(3)
Copy_ppt(15, 5)  # 田字格:5页
Copy_ppt(10, 1)  # 附录1:一页
Copy_ppt(11, 1)  # 附录2:一页

# pptSel.Save()  # 保存
pptSel.SaveAs(os.getcwd() + "\\" + "笔记本_增加页面temp.pptx")  # 另存为
pptSel.Close()  # 关闭 PowerPoint 文档
ppt.Quit()  # 关闭 office


Del_page()

# 生成文件为：笔记本_增加页面.pptx


# 删除不要的页面
# def del_slide(prs,index):
#     slides = list(prs.slides._sldIdLst)
#     prs.slides._sldIdLst.remove(slides[index])
#
# def Del_page():
#     # 打开ppt
#     ppt = Presentation('2笔记本_增加页面.pptx')
#
#     # 获取所有页
#     slides = ppt.slides
#     number_pages = len(slides)
#     print("删除前ppt一共",number_pages,"页面")
#
#     # 设置需要删除的页面数量
#     delPageNums = 12
#     # 进行删除操作（每次都删除第一张ppt）
#     for index in range(delPageNums):
#         del_slide(ppt,4)
#
#     # 再次获取所有页
#     slides = ppt.slides
#     number_pages = len(slides)
#     print("删除后ppt一共",number_pages,"页面")
#
#     ppt.save('2笔记本_增加页面.pptx')
#     print('删除完毕,生成新文件')


