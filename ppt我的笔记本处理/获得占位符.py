from pptx import Presentation
文件 = Presentation('笔记本 - 副本.pptx')
for index in range(0,len(文件.slides)):
    幻灯片=文件.slides[index]  #第4张幻灯片
# for shape in 幻灯片.shapes:
#     print(f'索引{shape},名称{shape.name},类型{shape.shape_type}')
    for 占位符 in 幻灯片.placeholders:
        信息 = 占位符.placeholder_format
    # 获得占位符信息
        print(f'幻灯片第{index}页(0为首页):索引{信息.idx},名称{占位符.name},类型{信息.type}')

    # 在占位符位置输出占位符索引号信息
        占位符.text=f'{信息.idx}'
    # 一个占位符输出格式信息
    # 幻灯片.placeholders[10].text='88'
文件.save("占位符.pptx")


