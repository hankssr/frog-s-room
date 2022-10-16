from pptx import Presentation
import pandas as pd
文件 = Presentation('笔记本_增加页面.pptx')
幻灯片=文件.slides[21]
p星期一=幻灯片.placeholders[10]
p星期二=幻灯片.placeholders[11]
p星期三=幻灯片.placeholders[12]
p星期四=幻灯片.placeholders[13]
p星期五=幻灯片.placeholders[14]
p星期六=幻灯片.placeholders[15]
p星期天=幻灯片.placeholders[16]

q星期一=幻灯片.placeholders[17]
q星期二=幻灯片.placeholders[18]
q星期三=幻灯片.placeholders[19]
q星期四=幻灯片.placeholders[20]
q星期五=幻灯片.placeholders[21]
q星期六=幻灯片.placeholders[22]
q星期天=幻灯片.placeholders[23]

m星期一=幻灯片.placeholders[24]
m星期二=幻灯片.placeholders[25]
m星期三=幻灯片.placeholders[26]
m星期四=幻灯片.placeholders[27]
m星期五=幻灯片.placeholders[28]
m星期六=幻灯片.placeholders[29]
m星期天=幻灯片.placeholders[30]

n星期一=幻灯片.placeholders[31]
n星期二=幻灯片.placeholders[32]
n星期三=幻灯片.placeholders[33]
n星期四=幻灯片.placeholders[34]
n星期五=幻灯片.placeholders[35]
n星期六=幻灯片.placeholders[36]
n星期天=幻灯片.placeholders[37]

o星期一=幻灯片.placeholders[38]
o星期二=幻灯片.placeholders[39]
o星期三=幻灯片.placeholders[40]
o星期四=幻灯片.placeholders[41]
o星期五=幻灯片.placeholders[42]
o星期六=幻灯片.placeholders[43]
o星期天=幻灯片.placeholders[44]

数据=pd.read_excel("日历.xlsx")
print(数据)
p星期一 = str(数据.iloc(0, 1))
p星期二 = str(数据.iloc(0, 2))
p星期三 = str(数据.iloc(0, 3))
p星期四 = str(数据.iloc(0, 4))
p星期五 = str(数据.iloc(0, 5))
p星期六 = str(数据.iloc(0, 6))
p星期天 = str(数据.iloc(0, 7))
q星期一 = str(数据.iloc(1, 1))
q星期二 = str(数据.iloc(1, 2))
q星期三 = str(数据.iloc(1, 3))
q星期四 = str(数据.iloc(1, 4))
q星期五 = str(数据.iloc(1, 5))
q星期六 = str(数据.iloc(1, 6))
q星期天 = str(数据.iloc(1, 7))
m星期一 = str(数据.iloc(2, 1))
m星期二 = str(数据.iloc(2, 2))
m星期三 = str(数据.iloc(2, 3))
m星期四 = str(数据.iloc(2, 4))
m星期五 = str(数据.iloc(2, 5))
m星期六 = str(数据.iloc(2, 6))
m星期天 = str(数据.iloc(2, 7))
n星期一 = str(数据.iloc(3, 1))
n星期二 = str(数据.iloc(3, 2))
n星期三 = str(数据.iloc(3, 3))
n星期四 = str(数据.iloc(3, 4))
n星期五 = str(数据.iloc(3, 5))
n星期六 = str(数据.iloc(3, 6))
n星期天 = str(数据.iloc(3, 7))
o星期一 = str(数据.iloc(4, 1))
o星期二 = str(数据.iloc(4, 2))
o星期三 = str(数据.iloc(4, 3))
o星期四 = str(数据.iloc(4, 4))
o星期五 = str(数据.iloc(4, 5))
o星期六 = str(数据.iloc(4, 6))
o星期天 = str(数据.iloc(4, 7))

文件.save("temp.pptx")
